"""Generate the binary sample templates (docx / pptx / xlsx).

These are plain Office files with Jinja2 ``{{ placeholders }}`` typed into them —
exactly what a caller would upload. Run once to populate the samples/ folder:

    python samples/make_samples.py

The text/HTML samples (letter.txt, report.html) are checked in directly.
"""
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from openpyxl.styles import Font
from pptx import Presentation

HERE = Path(__file__).resolve().parent


def make_docx() -> None:
    doc = Document()
    doc.add_heading("{{ company_name }}", level=0)
    doc.add_paragraph("Invoice {{ invoice_number }} — {{ date }}")
    doc.add_paragraph("Billed to: {{ customer_name }}")
    doc.add_paragraph("")
    doc.add_paragraph("Items:")
    # Single-paragraph loop keeps run formatting predictable.
    doc.add_paragraph(
        "{% for item in items %}- {{ item.name }}: {{ item.amount }}\n{% endfor %}"
    )
    doc.add_paragraph("Total due: {{ total }}")
    doc.save(HERE / "invoice.docx")


def make_pptx() -> None:
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "{{ company_name }}"
    title_slide.placeholders[1].text = "Invoice {{ invoice_number }} for {{ customer_name }}"

    body_slide = prs.slides.add_slide(prs.slide_layouts[1])
    body_slide.shapes.title.text = "Summary — {{ date }}"
    body = body_slide.placeholders[1].text_frame
    body.text = "Customer: {{ customer_name }}"
    body.add_paragraph().text = "Invoice: {{ invoice_number }}"
    body.add_paragraph().text = "Total due: {{ total }}"
    prs.save(HERE / "deck.pptx")


def make_xlsx() -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Invoice"
    ws["A1"] = "{{ company_name }}"
    ws["A1"].font = Font(bold=True, size=14)
    ws["A2"] = "Invoice"
    ws["B2"] = "{{ invoice_number }}"
    ws["A3"] = "Date"
    ws["B3"] = "{{ date }}"
    ws["A4"] = "Customer"
    ws["B4"] = "{{ customer_name }}"
    ws["A6"] = "Total due"
    ws["B6"] = "{{ total }}"
    wb.save(HERE / "sheet.xlsx")


if __name__ == "__main__":
    make_docx()
    make_pptx()
    make_xlsx()
    print(f"Wrote invoice.docx, deck.pptx, sheet.xlsx to {HERE}")
