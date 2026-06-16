"""PPTX generator.

Walks every slide (shape text frames and table cells) of an uploaded .pptx
template and renders any paragraph containing a Jinja2 marker against the
supplied content.

v1 limitation: rendering happens at the paragraph level. When a marker spans
multiple text runs with different formatting, the rendered text inherits the
first run's formatting and the remaining runs are emptied. Keep each placeholder
within a single run/format for best fidelity.
"""
from __future__ import annotations

from io import BytesIO

from jinja2 import BaseLoader, Environment, StrictUndefined
from pptx import Presentation

_env = Environment(loader=BaseLoader(), undefined=StrictUndefined, autoescape=False)


def _has_marker(value: str) -> bool:
    return "{{" in value or "{%" in value


def _process_text_frame(text_frame, content: dict) -> None:
    for para in text_frame.paragraphs:
        full_text = "".join(run.text for run in para.runs)
        if not _has_marker(full_text):
            continue
        rendered = _env.from_string(full_text).render(**content)
        if para.runs:
            para.runs[0].text = rendered
            for run in para.runs[1:]:
                run.text = ""
        else:
            para.add_run().text = rendered


def generate(template_bytes: bytes, content: dict) -> bytes:
    prs = Presentation(BytesIO(template_bytes))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                _process_text_frame(shape.text_frame, content)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        _process_text_frame(cell.text_frame, content)
    out = BytesIO()
    prs.save(out)
    return out.getvalue()
