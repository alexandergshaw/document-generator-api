"""End-to-end tests for the document generator API.

Templates are built in-memory so the suite is self-contained (no dependency on
the samples/ folder having been generated).
"""
import io
import json

import pytest
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation

import app as flask_app

CONTENT = {
    "company_name": "Acme Corp",
    "customer_name": "Ada Lovelace",
    "invoice_number": "A-1042",
    "date": "2026-06-15",
    "items": [
        {"name": "Widget", "amount": "$10.00"},
        {"name": "Gadget", "amount": "$25.50"},
    ],
    "total": "$35.50",
}


@pytest.fixture
def client():
    flask_app.app.config.update(TESTING=True)
    return flask_app.app.test_client()


@pytest.fixture(autouse=True)
def reset_config():
    """Keep optional auth/CORS off by default and restore after each test."""
    app = flask_app.app
    saved = {k: app.config.get(k) for k in
             ("API_KEY", "ALLOWED_ORIGINS", "MAX_CONTENT_LENGTH")}
    app.config["API_KEY"] = None
    app.config["ALLOWED_ORIGINS"] = []
    yield
    app.config.update(saved)


# --- in-memory template builders ------------------------------------------

def _docx_template() -> bytes:
    doc = Document()
    doc.add_paragraph("Invoice {{ invoice_number }} for {{ customer_name }}")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _pptx_template() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = "{{ company_name }}"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _xlsx_template() -> bytes:
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "{{ company_name }}"
    ws["A2"] = "{{ invoice_number }}"
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# --- format discovery ------------------------------------------------------

def test_formats_endpoint(client):
    res = client.get("/api/formats")
    assert res.status_code == 200
    names = {f["name"] for f in res.get_json()["formats"]}
    assert {"txt", "md", "html", "csv", "docx", "pptx", "xlsx", "pdf"} <= names


# --- text-family formats ---------------------------------------------------

def test_generate_txt(client):
    res = client.post("/api/generate", data={
        "document_type": "txt",
        "content": json.dumps(CONTENT),
        "template_text": "Hello {{ customer_name }} ({{ invoice_number }})",
    })
    assert res.status_code == 200
    assert res.mimetype == "text/plain"
    assert b"Ada Lovelace" in res.data
    assert b"A-1042" in res.data


def test_generate_html(client):
    res = client.post("/api/generate", data={
        "document_type": "html",
        "content": json.dumps(CONTENT),
        "template_text": "<h1>{{ company_name }}</h1>",
    })
    assert res.status_code == 200
    assert res.mimetype == "text/html"
    assert b"Acme Corp" in res.data


def test_generate_pdf(client):
    res = client.post("/api/generate", data={
        "document_type": "pdf",
        "content": json.dumps(CONTENT),
        "template_text": "<h1>{{ company_name }}</h1><p>{{ invoice_number }}</p>",
    })
    assert res.status_code == 200
    assert res.mimetype == "application/pdf"
    assert res.data[:4] == b"%PDF"


# --- office formats --------------------------------------------------------

def test_generate_docx(client):
    res = client.post("/api/generate", data={
        "document_type": "docx",
        "content": json.dumps(CONTENT),
        "template": (io.BytesIO(_docx_template()), "t.docx"),
    }, content_type="multipart/form-data")
    assert res.status_code == 200
    assert res.data[:2] == b"PK"  # .docx is a zip container
    doc = Document(io.BytesIO(res.data))
    text = "\n".join(p.text for p in doc.paragraphs)
    assert "A-1042" in text
    assert "Ada Lovelace" in text


def test_generate_xlsx(client):
    res = client.post("/api/generate", data={
        "document_type": "xlsx",
        "content": json.dumps(CONTENT),
        "template": (io.BytesIO(_xlsx_template()), "t.xlsx"),
    }, content_type="multipart/form-data")
    assert res.status_code == 200
    ws = load_workbook(io.BytesIO(res.data)).active
    assert ws["A1"].value == "Acme Corp"
    assert ws["A2"].value == "A-1042"


def test_generate_pptx(client):
    res = client.post("/api/generate", data={
        "document_type": "pptx",
        "content": json.dumps(CONTENT),
        "template": (io.BytesIO(_pptx_template()), "t.pptx"),
    }, content_type="multipart/form-data")
    assert res.status_code == 200
    prs = Presentation(io.BytesIO(res.data))
    texts = [
        s.text_frame.text
        for slide in prs.slides
        for s in slide.shapes
        if s.has_text_frame
    ]
    assert any("Acme Corp" in t for t in texts)


# --- error handling --------------------------------------------------------

def test_unsupported_type(client):
    res = client.post("/api/generate", data={
        "document_type": "xyz", "content": "{}", "template_text": "x",
    })
    assert res.status_code == 400
    assert "Unsupported" in res.get_json()["error"]


def test_bad_json(client):
    res = client.post("/api/generate", data={
        "document_type": "txt", "content": "{not json", "template_text": "x",
    })
    assert res.status_code == 400
    assert "valid JSON" in res.get_json()["error"]


def test_missing_required_template(client):
    res = client.post("/api/generate", data={
        "document_type": "docx", "content": "{}",
    })
    assert res.status_code == 400
    assert "requires an uploaded template" in res.get_json()["error"]


def test_missing_content_key_is_error(client):
    # StrictUndefined: a placeholder with no matching key surfaces as an error.
    res = client.post("/api/generate", data={
        "document_type": "txt",
        "content": "{}",
        "template_text": "Hi {{ missing_key }}",
    })
    assert res.status_code == 500
    assert "Generation failed" in res.get_json()["error"]


# --- versioning ------------------------------------------------------------

def test_health_returns_version(client):
    res = client.get("/api/health")
    assert res.status_code == 200
    body = res.get_json()
    assert body["status"] == "ok"
    assert isinstance(body["version"], str) and body["version"]


def test_formats_includes_version(client):
    res = client.get("/api/formats")
    assert isinstance(res.get_json().get("version"), str)


# --- uniform error envelope + machine codes --------------------------------

def test_code_unsupported_type(client):
    res = client.post("/api/generate", data={
        "document_type": "xyz", "template_text": "x", "content": "{}"})
    assert res.status_code == 400
    assert res.get_json()["code"] == "unsupported_type"


def test_code_invalid_json(client):
    res = client.post("/api/generate", data={
        "document_type": "txt", "content": "{not json", "template_text": "x"})
    assert res.status_code == 400
    assert res.get_json()["code"] == "invalid_json"


def test_code_content_not_object(client):
    res = client.post("/api/generate", data={
        "document_type": "txt", "content": "[1, 2, 3]", "template_text": "x"})
    assert res.status_code == 400
    assert res.get_json()["code"] == "content_not_object"


def test_code_template_required(client):
    res = client.post("/api/generate", data={
        "document_type": "docx", "content": "{}"})
    assert res.status_code == 400
    assert res.get_json()["code"] == "template_required"


def test_code_no_template(client):
    res = client.post("/api/generate", data={
        "document_type": "txt", "content": "{}"})
    assert res.status_code == 400
    assert res.get_json()["code"] == "no_template"


def test_code_render_failed(client):
    res = client.post("/api/generate", data={
        "document_type": "txt", "content": "{}", "template_text": "Hi {{ nope }}"})
    assert res.status_code == 500
    assert res.get_json()["code"] == "render_failed"


def test_413_returns_json_too_large(client):
    flask_app.app.config["MAX_CONTENT_LENGTH"] = 64
    res = client.post("/api/generate", data={
        "document_type": "txt", "content": "{}", "template_text": "x" * 500})
    assert res.status_code == 413
    assert res.mimetype == "application/json"
    assert res.get_json()["code"] == "too_large"


def test_unknown_route_returns_json(client):
    res = client.get("/api/does-not-exist")
    assert res.status_code == 404
    assert res.mimetype == "application/json"
    assert "code" in res.get_json()


# --- configurable strictness ----------------------------------------------

def test_strict_false_blanks_txt(client):
    res = client.post("/api/generate", data={
        "document_type": "txt", "content": "{}",
        "template_text": "A{{ missing }}B", "strict": "false"})
    assert res.status_code == 200
    assert res.data == b"AB"


def test_strict_false_blanks_pdf(client):
    res = client.post("/api/generate", data={
        "document_type": "pdf", "content": "{}",
        "template_text": "<p>{{ missing }}</p>", "strict": "false"})
    assert res.status_code == 200
    assert res.data[:4] == b"%PDF"


def test_docx_default_blanks_missing(client):
    # docx default (no strict field) renders a missing key blank -> success.
    res = client.post("/api/generate", data={
        "document_type": "docx", "content": "{}",
        "template": (io.BytesIO(_docx_template()), "t.docx"),
    }, content_type="multipart/form-data")
    assert res.status_code == 200


def test_strict_true_errors_docx(client):
    res = client.post("/api/generate", data={
        "document_type": "docx", "content": "{}", "strict": "true",
        "template": (io.BytesIO(_docx_template()), "t.docx"),
    }, content_type="multipart/form-data")
    assert res.status_code == 500
    assert res.get_json()["code"] == "render_failed"


# --- filename sanitization -------------------------------------------------

def test_filename_traversal_neutralized(client):
    res = client.post("/api/generate", data={
        "document_type": "txt", "content": "{}", "template_text": "hi",
        "filename": "../../etc/passwd"})
    assert res.status_code == 200
    disposition = res.headers["Content-Disposition"]
    assert "passwd.txt" in disposition
    assert ".." not in disposition
    assert "/" not in disposition.split("filename=")[-1]


# --- optional auth ---------------------------------------------------------

def test_auth_rejects_bad_key(client):
    flask_app.app.config["API_KEY"] = "secret"
    res = client.post("/api/generate", data={
        "document_type": "txt", "content": "{}", "template_text": "hi"})
    assert res.status_code == 401
    assert res.get_json()["code"] == "unauthorized"


def test_auth_accepts_good_key(client):
    flask_app.app.config["API_KEY"] = "secret"
    res = client.post("/api/generate", data={
        "document_type": "txt", "content": "{}", "template_text": "hi"},
        headers={"X-API-Key": "secret"})
    assert res.status_code == 200


def test_no_auth_required_when_unset(client):
    res = client.post("/api/generate", data={
        "document_type": "txt", "content": "{}", "template_text": "hi"})
    assert res.status_code == 200


# --- optional CORS ---------------------------------------------------------

def test_cors_headers_when_configured(client):
    flask_app.app.config["ALLOWED_ORIGINS"] = ["https://app.example"]
    res = client.get("/api/formats", headers={"Origin": "https://app.example"})
    assert res.headers.get("Access-Control-Allow-Origin") == "https://app.example"
    assert "X-API-Key" in res.headers.get("Access-Control-Allow-Headers", "")


def test_no_cors_headers_by_default(client):
    res = client.get("/api/formats", headers={"Origin": "https://app.example"})
    assert "Access-Control-Allow-Origin" not in res.headers


def test_cors_ignores_unlisted_origin(client):
    flask_app.app.config["ALLOWED_ORIGINS"] = ["https://app.example"]
    res = client.get("/api/formats", headers={"Origin": "https://evil.example"})
    assert "Access-Control-Allow-Origin" not in res.headers
